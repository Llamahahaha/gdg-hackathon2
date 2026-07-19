from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    title_only_layout = prs.slide_layouts[5]
    
    # 1. Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "FieldTheory AI (VisionPlay)"
    subtitle.text = "Spatio-Temporal Graph Intelligence Engine\nSoftware Engineering & Project Management Presentation\nTeam Size: 3 Engineers"
    
    # 2. Introduction & Scope
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Introduction & Scope of Software"
    tf = body_shape.text_frame
    tf.text = "FieldTheory AI is an elite-tier tactical intelligence platform for sports."
    p = tf.add_paragraph()
    p.text = "Scope & Economic Aspects: Moves beyond event-based stats (e.g., pass completion) to structural health and graph theory, providing venture-scale value to sports analysts."
    p = tf.add_paragraph()
    p.text = "Core Functionality: Translates raw video footage into a temporal weighted graph to predict tactical collapses."

    # 3. Software Life Cycle Model
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Software Life Cycle Model (SDLC)"
    tf = body_shape.text_frame
    tf.text = "Model Used: Rapid Prototyping & Agile/Iterative Model"
    p = tf.add_paragraph()
    p.text = "Why?: Built under strict hackathon time constraints, requiring rapid iteration between the CV pipeline and the UI."
    p = tf.add_paragraph()
    p.text = "Extreme Programming (XP) Elements: Continuous integration (Next.js + FastAPI), pair programming, and rapid refactoring (e.g., swapping Turbopack for Webpack on the fly)."

    # 4. Team Organization & Management
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Team Organization & Management"
    tf = body_shape.text_frame
    tf.text = "Democratic Team Approach (Team of 3 Engineers)"
    p = tf.add_paragraph()
    p.text = "Collaboration: Equal weight in design decisions, with specialized roles (Frontend UI, Graph Math Engine, CV Pipeline)."
    p = tf.add_paragraph()
    p.text = "Synchronize & Stabilize: Daily syncing of the WebSocket endpoints between the Python backend and React frontend."
    
    # 5. Requirement & Specification Phase
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Requirement Elicitation & Analysis"
    tf = body_shape.text_frame
    tf.text = "Client/User Needs: Coaches need visual, real-time tactical intelligence, not spreadsheets."
    p = tf.add_paragraph()
    p.text = "Use Case Modeling:"
    p.level = 1
    p.text = "Actor: Football Analyst / Coach"
    p.level = 1
    p.text = "Action: Uploads footage -> Views Graph Entropy -> Simulates Lynchpin removal."
    p = tf.add_paragraph()
    p.text = "Non-Functional Requirements: High-performance streaming (low latency WebSockets), visual cinematic quality (Framer Motion)."

    # 6. Object-Oriented Design & Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "OO Design, Cohesion & Coupling"
    tf = body_shape.text_frame
    tf.text = "Architecture: Client-Server Model (Next.js + FastAPI + Supabase)"
    p = tf.add_paragraph()
    p.text = "High Cohesion: The CV Engine (YOLO) and Graph Engine (NetworkX) are highly cohesive, handling distinct logical tasks."
    p = tf.add_paragraph()
    p.text = "Low Coupling: The frontend UI is fully decoupled from the mathematical engine, communicating strictly via WebSocket JSON payloads."
    p = tf.add_paragraph()
    p.text = "Data Encapsulation: Tactical context and authentication states are encapsulated using React Context (e.g., AuthContext, TacticalContext)."

    # 7. Algorithms & Core Logic (Implementation)
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Algorithms & Implementation"
    tf = body_shape.text_frame
    tf.text = "Graph Laplacian & Spectral Entropy: Evaluates overall formation stability (Fiedler value)."
    p = tf.add_paragraph()
    p.text = "Tarjan’s Algorithm: Identifies Articulation Points (Lynchpin players)."
    p = tf.add_paragraph()
    p.text = "Floyd-Warshall Algorithm: Computes the Team Diameter to measure compactness."

    # 8. Software Testing Strategies
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Software Testing & SQA"
    tf = body_shape.text_frame
    tf.text = "Testing Strategy: Incremental and Integration Testing"
    p = tf.add_paragraph()
    p.text = "Integration Testing: Verifying the WebSocket streaming between Uvicorn/FastAPI and the Next.js React-Flow graphs."
    p = tf.add_paragraph()
    p.text = "Black Box Testing: Simulating user clicks on 'Neutralize Player' to ensure the graph fractures correctly without crashing."
    p = tf.add_paragraph()
    p.text = "Exception Handling: Safely catching Supabase database write errors and Turbopack compiler panics."

    # 9. Project Management & Risk Management
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Risk Management & PM"
    tf = body_shape.text_frame
    tf.text = "Requirement Change Management: Pivoted from a generic analytics dashboard to a 'Spatio-Temporal Graph Engine' for competitive advantage."
    p = tf.add_paragraph()
    p.text = "Risk Identification & Mitigation:"
    p.level = 1
    p.text = "Risk: Firebase API keys were suspended mid-development."
    p.level = 1
    p.text = "Mitigation: Implemented a robust Mock Auth bypass to ensure the presentation/demo remained flawless."
    p = tf.add_paragraph()
    p.text = "Configuration Management: Managed `.env` variables for Supabase and local WebSocket routing."

    # 10. Conclusion & Closure
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Project Closure & Future Scope"
    tf = body_shape.text_frame
    tf.text = "Current Status: Live Engine and Dashboard are fully operational."
    p = tf.add_paragraph()
    p.text = "Future Scope (Maintenance Phase):"
    p.level = 1
    p.text = "Expand to live edge-stream processing."
    p.level = 1
    p.text = "PDF Tactical Audits and full interactive Sandbox."
    p = tf.add_paragraph()
    p.text = "Closure Analysis: The project successfully integrates advanced mathematical theory with a high-performance web interface, satisfying the SEPM paradigms of robust architecture, iterative prototyping, and decoupled design."

    prs.save('SEPM_Project_Presentation.pptx')

if __name__ == '__main__':
    create_presentation()
